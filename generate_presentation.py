from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from novaclose_analysis import (
    build_scenario_actions,
    calculate_kpis,
    calculate_readiness_score,
    load_data,
    priority_engine,
    simulate_close_scenario,
)


ROOT = Path(__file__).resolve().parent
DATASET = ROOT / "AI_Finance_Hackathon_Month_End_Dataset.xlsx"
OUTPUT = ROOT / "NovaClose_App_Presentation.pptx"

NAVY = RGBColor(20, 54, 66)
TEAL = RGBColor(42, 157, 143)
ORANGE = RGBColor(215, 122, 45)
RED = RGBColor(199, 92, 92)
PURPLE = RGBColor(115, 83, 229)
INK = RGBColor(15, 31, 44)
MUTED = RGBColor(91, 107, 120)
SURFACE = RGBColor(245, 249, 252)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(214, 223, 229)


def fmt_int(value):
    return f"{int(value):,}"


def fmt_num(value, decimals=1):
    return f"{value:,.{decimals}f}"


def add_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = SURFACE


def add_title(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.45), Inches(8.5), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.clear()
    p = title_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = INK

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.62), Inches(1.0), Inches(10.8), Inches(0.45))
        sub_frame = sub_box.text_frame
        sub_frame.clear()
        p = sub_frame.paragraphs[0]
        run = p.add_run()
        run.text = subtitle
        run.font.size = Pt(11)
        run.font.color.rgb = MUTED


def add_footer(slide, text):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(7.0), Inches(11.6), Inches(0.3))
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(9)
    run.font.color.rgb = MUTED


def add_card(slide, left, top, width, height, title, value=None, body=None, accent=NAVY):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LINE
    shape.line.width = Pt(1)

    title_box = slide.shapes.add_textbox(left + Inches(0.16), top + Inches(0.1), width - Inches(0.3), Inches(0.25))
    title_frame = title_box.text_frame
    title_frame.clear()
    p = title_frame.paragraphs[0]
    run = p.add_run()
    run.text = title.upper()
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = MUTED

    if value is not None:
        value_box = slide.shapes.add_textbox(left + Inches(0.16), top + Inches(0.42), width - Inches(0.3), Inches(0.45))
        value_frame = value_box.text_frame
        value_frame.clear()
        p = value_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(value)
        run.font.size = Pt(24)
        run.font.bold = True
        run.font.color.rgb = accent

    if body:
        body_box = slide.shapes.add_textbox(left + Inches(0.16), top + Inches(0.92), width - Inches(0.3), height - Inches(1.05))
        body_frame = body_box.text_frame
        body_frame.clear()
        p = body_frame.paragraphs[0]
        run = p.add_run()
        run.text = body
        run.font.size = Pt(10)
        run.font.color.rgb = MUTED


def add_bullets(slide, left, top, width, height, lines, title=None):
    if title:
        label = slide.shapes.add_textbox(left, top, width, Inches(0.3))
        frame = label.text_frame
        frame.clear()
        p = frame.paragraphs[0]
        run = p.add_run()
        run.text = title.upper()
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.color.rgb = MUTED
        top += Inches(0.3)
        height -= Inches(0.3)

    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    for index, line in enumerate(lines):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(15)
        p.font.color.rgb = INK
        p.space_after = Pt(8)


def build_deck():
    data = load_data(str(DATASET))
    kpis = calculate_kpis(data)
    score = calculate_readiness_score(kpis)
    priorities = priority_engine(kpis)
    actions = build_scenario_actions(kpis)
    scenario = simulate_close_scenario(kpis, [action["id"] for action in actions])

    summary = kpis["summary"]
    scenario_summary = scenario["kpis"]["summary"]
    posting_summary = scenario["posting_simulator"]["summary"]
    period = summary["accounting_period_label"]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    banner = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(0.5), Inches(12.2), Inches(1.25))
    banner.fill.solid()
    banner.fill.fore_color.rgb = WHITE
    banner.line.color.rgb = LINE
    title = slide.shapes.add_textbox(Inches(0.85), Inches(0.82), Inches(7.8), Inches(0.5))
    tf = title.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "NovaClose"
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = INK
    sub = slide.shapes.add_textbox(Inches(0.88), Inches(1.2), Inches(9.0), Inches(0.35))
    tf = sub.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"AI-powered month-end close command centre | NovaTech | {period}"
    run.font.size = Pt(13)
    run.font.color.rgb = MUTED
    add_card(slide, Inches(0.7), Inches(2.0), Inches(3.75), Inches(1.45), "Close Readiness", f"{score['readiness_score']}/100", "Baseline command-centre score before automation.", RED)
    add_card(slide, Inches(4.8), Inches(2.0), Inches(3.75), Inches(1.45), "Projected Close", f"{score['predicted_close_days']} days", "Current dashboard forecast versus the 4-day target.", ORANGE)
    add_card(slide, Inches(8.9), Inches(2.0), Inches(3.75), Inches(1.45), "After Automation", f"{scenario['score']['predicted_close_days']} days", "Forecast after the full automation bundle is executed.", TEAL)
    add_bullets(
        slide,
        Inches(0.82),
        Inches(4.0),
        Inches(11.7),
        Inches(2.2),
        [
            "Transforms month-end from reactive exception hunting into a guided, auditable operating system.",
            "Combines Command Centre visibility, specialist agents, before/after scenario testing, and AI copilot support.",
        ],
        title="Why this app exists",
    )
    add_footer(slide, "NovaClose | App presentation")

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "What the app is solving", "The delay is caused by exception handling and approval friction, not transaction processing.")
    add_card(slide, Inches(0.65), Inches(1.45), Inches(2.0), Inches(1.35), "Pending GL", fmt_int(summary["pending_gl"]), "Journals still in review or approval.", RED)
    add_card(slide, Inches(2.85), Inches(1.45), Inches(2.0), Inches(1.35), "Bank Exceptions", fmt_int(summary["bank_exceptions"]), "Cash items still open in reconciliation.", ORANGE)
    add_card(slide, Inches(5.05), Inches(1.45), Inches(2.0), Inches(1.35), "IC Exceptions", fmt_int(summary["ic_exceptions"]), "Cross-entity differences still unresolved.", RED)
    add_card(slide, Inches(7.25), Inches(1.45), Inches(2.0), Inches(1.35), "AP Issues", fmt_int(summary["ap_3way_exceptions"]), "3-way match exceptions still open.", ORANGE)
    add_card(slide, Inches(9.45), Inches(1.45), Inches(2.0), Inches(1.35), "Blocked Tasks", fmt_int(summary["blocked_tasks"]), "Checklist dependencies blocking close.", RED)
    add_bullets(
        slide,
        Inches(0.75),
        Inches(3.15),
        Inches(5.6),
        Inches(3.2),
        [
            "171 GL approvals and reviews are the biggest operational blocker.",
            "Cash, intercompany, and AP exceptions create downstream reporting delays.",
            "Checklist handoffs break because issues are found late, not because finance lacks activity.",
        ],
        title="Root cause",
    )
    add_bullets(
        slide,
        Inches(6.8),
        Inches(3.15),
        Inches(5.7),
        Inches(3.2),
        [
            "CFO needs one place to see readiness, top blockers, and the path below 4 days.",
            "Ops team needs to know exactly what to clear next and which actions can be automated safely.",
            "Audit needs transparency on what moved automatically and what still needs sign-off.",
        ],
        title="What the app delivers",
    )
    add_footer(slide, "NovaClose | Problem framing")

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Before and after the automation bundle", "The app shows the effect of running the operational bundle, not just descriptive insight.")
    add_card(slide, Inches(0.7), Inches(1.5), Inches(2.0), Inches(1.45), "Readiness", f"{score['readiness_score']}/100", "Before automation", RED)
    add_card(slide, Inches(2.9), Inches(1.5), Inches(2.0), Inches(1.45), "Readiness", f"{scenario['score']['readiness_score']}/100", "After automation", TEAL)
    add_card(slide, Inches(5.1), Inches(1.5), Inches(2.0), Inches(1.45), "Continuous Forecast", f"{fmt_num(score['continuous_close_days'])} days", "Before automation", ORANGE)
    add_card(slide, Inches(7.3), Inches(1.5), Inches(2.0), Inches(1.45), "Continuous Forecast", f"{fmt_num(scenario['score']['continuous_close_days'])} days", "After automation", TEAL)
    add_card(slide, Inches(9.5), Inches(1.5), Inches(2.0), Inches(1.45), "Dashboard Forecast", f"{scenario['score']['predicted_close_days']} days", "After automation", TEAL)
    compare_rows = [
        ("Pending GL", summary["pending_gl"], scenario_summary["pending_gl"]),
        ("Bank Exceptions", summary["bank_exceptions"], scenario_summary["bank_exceptions"]),
        ("IC Exceptions", summary["ic_exceptions"], scenario_summary["ic_exceptions"]),
        ("AP Exceptions", summary["ap_3way_exceptions"], scenario_summary["ap_3way_exceptions"]),
        ("Blocked Tasks", summary["blocked_tasks"], scenario_summary["blocked_tasks"]),
    ]
    table = slide.shapes.add_table(len(compare_rows) + 1, 3, Inches(0.8), Inches(3.35), Inches(5.9), Inches(2.5)).table
    table.columns[0].width = Inches(2.8)
    table.columns[1].width = Inches(1.4)
    table.columns[2].width = Inches(1.7)
    headers = ["Metric", "Before", "After"]
    for col, text in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE
        cell.text_frame.paragraphs[0].font.size = Pt(11)
    for row_index, (label, before, after) in enumerate(compare_rows, start=1):
        table.cell(row_index, 0).text = label
        table.cell(row_index, 1).text = fmt_int(before)
        table.cell(row_index, 2).text = fmt_int(after)
        for col in range(3):
            cell = table.cell(row_index, col)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.color.rgb = INK
    add_bullets(
        slide,
        Inches(7.35),
        Inches(3.35),
        Inches(5.15),
        Inches(2.7),
        [
            f"Readiness improves by {scenario['score']['readiness_score'] - score['readiness_score']} points.",
            f"Continuous forecast improves from {fmt_num(score['continuous_close_days'])} to {fmt_num(scenario['score']['continuous_close_days'])} days.",
            "The biggest improvement comes from GL approval routing, checklist unblocking, and exception auto-clear.",
        ],
        title="Executive takeaway",
    )
    add_footer(slide, "NovaClose | Before vs after")

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "App structure", "The app is cleaner for demo use; explanatory narrative is carried in this deck instead.")
    modules = [
        ("Command Centre", "Readiness, before/after close signal, month-end steps, and action queue", NAVY),
        ("Close Tracker", "Checklist blockers, dependency handoffs, and critical-path work", ORANGE),
        ("Agents", "GL, Bank, IC, Journal, Audit, and Variance Analysis specialist views", TEAL),
        ("Risks", "Eight-sheet risk scan plus legal-entity risk comparison", RED),
        ("Scenarios", "Automation bundle impact, posting outcomes, and before/after evidence", PURPLE),
        ("AI Copilot", "Finance Q&A routed through the local close-intelligence engine", NAVY),
    ]
    positions = [
        (Inches(0.75), Inches(1.55)),
        (Inches(4.45), Inches(1.55)),
        (Inches(8.15), Inches(1.55)),
        (Inches(0.75), Inches(4.0)),
        (Inches(4.45), Inches(4.0)),
        (Inches(8.15), Inches(4.0)),
    ]
    for (title, body, accent), (left, top) in zip(modules, positions):
        add_card(slide, left, top, Inches(3.1), Inches(1.95), title, None, body, accent)
    add_footer(slide, "NovaClose | Product walkthrough")

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Agent suite", "Each agent owns a specific part of the close and feeds a common readiness and posting model.")
    agent_rows = [
        ("GL Approval Agent", "Triages pending journals into straight-through, manager, controller, and CFO lanes."),
        ("Bank Reconciliation Agent", "Routes cash exceptions, drafts residual journals, and isolates statement breaks."),
        ("IC Reconciliation Agent", "Auto-matches counterparties, resolves FX mismatches, and drafts eliminations."),
        ("Journal Entry Agent", "Builds ERP-ready accrual and reclass drafts with support and audit trails."),
        ("Audit & Compliance Agent", "Scores control completeness and separates ready-to-post from held entries."),
        ("Variance Analysis Agent", "Calculates MoM and YoY anomalies and drafts executive commentary."),
        ("Close Tracker", "Turns checklist dependencies into an unblock queue with critical-path visibility."),
    ]
    add_bullets(
        slide,
        Inches(0.8),
        Inches(1.5),
        Inches(12.0),
        Inches(5.5),
        [f"{name}: {copy}" for name, copy in agent_rows],
        title="What moved out of the UI and into the deck",
    )
    add_footer(slide, "NovaClose | Agent architecture")

    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Automation bundle and control transparency", "One controlled bundle runs in the app, with an audit trail on what changed and what still needs release.")
    bundle_lines = [f"{index + 1}. {action['title']}" for index, action in enumerate(actions)]
    add_bullets(slide, Inches(0.75), Inches(1.45), Inches(7.0), Inches(4.8), bundle_lines, title="Bundle actions")
    add_card(slide, Inches(8.3), Inches(1.55), Inches(2.0), Inches(1.4), "Auto-Post", fmt_int(posting_summary["auto_post"]), "Moves straight into ERP when policy allows.", TEAL)
    add_card(slide, Inches(10.45), Inches(1.55), Inches(2.0), Inches(1.4), "Ready to Post", fmt_int(posting_summary["ready_to_post"]), "Cleared but still awaiting release.", ORANGE)
    add_card(slide, Inches(8.3), Inches(3.2), Inches(2.0), Inches(1.4), "Manual Hold", fmt_int(posting_summary["manual_hold"]), "Still requires review or approval.", RED)
    add_card(slide, Inches(10.45), Inches(3.2), Inches(2.0), Inches(1.4), "Hours Saved", fmt_int(scenario["total_hours_saved"]), "Recovered from the automation bundle.", TEAL)
    add_bullets(
        slide,
        Inches(8.05),
        Inches(5.0),
        Inches(4.5),
        Inches(1.3),
        [
            "Run sheet logs timestamps, actions applied, hours saved, and posting state.",
            "Before/after risk cards and agent tables make the automation effect transparent.",
        ],
    )
    add_footer(slide, "NovaClose | Automation controls")

    # Slide 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "How the CFO and operations team use it", "The app is designed to answer three questions fast: Where are we? What moved? What do we do next?")
    add_bullets(
        slide,
        Inches(0.8),
        Inches(1.55),
        Inches(5.8),
        Inches(4.7),
        [
            "1. Load the workbook and validate the period.",
            "2. Review the Command Centre before/after readiness signal.",
            "3. Run the automation bundle from Month-End Steps.",
            "4. Work the top three priorities and owner queue.",
            "5. Use AI Copilot for blocker, entity, and variance questions.",
        ],
        title="Operator flow",
    )
    top_three = priorities[:3]
    add_bullets(
        slide,
        Inches(6.8),
        Inches(1.55),
        Inches(5.7),
        Inches(4.7),
        [f"{item['priority_item']} - {item['downstream_unlock']}" for item in top_three],
        title="Top focus areas from the current data",
    )
    add_footer(slide, "NovaClose | Operating model")

    # Slide 8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Business impact and demo message", "Use this slide as the closing narrative in the hackathon pitch.")
    add_card(slide, Inches(0.75), Inches(1.5), Inches(3.7), Inches(1.45), "Readiness Lift", f"+{scenario['score']['readiness_score'] - score['readiness_score']} pts", "68/100 to 92/100 after the bundle.", TEAL)
    add_card(slide, Inches(4.8), Inches(1.5), Inches(3.7), Inches(1.45), "Close Forecast", f"{fmt_num(score['continuous_close_days'])} -> {fmt_num(scenario['score']['continuous_close_days'])} days", "Continuous forecast after automation.", TEAL)
    add_card(slide, Inches(8.85), Inches(1.5), Inches(3.7), Inches(1.45), "Blockers Removed", "IC 2 -> 0 | Blocked 3 -> 0", "Exception and checklist drag removed from the critical path.", TEAL)
    add_bullets(
        slide,
        Inches(0.82),
        Inches(3.45),
        Inches(11.8),
        Inches(2.8),
        [
            "NovaClose gives the CFO a single view of readiness, forecast, and automation impact.",
            "It gives finance operations a clear next-action queue instead of a static exception report.",
            "It keeps control transparency visible by separating auto-post, ready-to-post, and manual-hold outcomes.",
            "The app is intentionally cleaner now; the deck carries the architecture, narrative, and agent story.",
        ],
        title="Pitch close",
    )
    add_footer(slide, "NovaClose | Hackathon close")

    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    path = build_deck()
    print(path)
